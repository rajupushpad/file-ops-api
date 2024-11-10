import io
from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework import status
from django.http import FileResponse
import os
from django.http import HttpResponse, JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.core.files.storage import default_storage

from pdf2docx import Converter
import PyPDF2;
from PyPDF2 import PdfReader, PdfWriter
from docx import Document

from pptx import Presentation
from pptx.util import Inches
import pdfplumber
# from openpyxl import Workbook
from PIL import Image

# not in use
import numpy as np

# from reportlab.lib.pagesizes import letter
# from reportlab.pdfgen import canvas

# import tabula
# import pandas as pd

# import tempfile

from pdf2image import convert_from_path

# import cairosvg
# import potrace
# import subprocess

from PIL import Image
import fitz
from io import BytesIO  # Import BytesIO for in-memory byte streams

# from fpdf import FPDF
# from docx import Document as DocxDocument

@api_view(['POST'])
def convert_pdf_to_docx(request):
    
    if 'file' not in request.FILES:
        return Response({'error': 'No file uploaded'}, status=status.HTTP_400_BAD_REQUEST)

    pdf_file = request.FILES['file']
    if not pdf_file.name.endswith('.pdf'):
        return Response({'error': 'Invalid file format'}, status=status.HTTP_400_BAD_REQUEST)

    try:
        
        # Save the file temporarily
        pdf_path = default_storage.save(pdf_file.name, pdf_file)
        full_pdf_path = os.path.join(default_storage.location, pdf_path)

        # Define output DOCX path
        docx_filename = pdf_file.name.replace('.pdf', '.docx')
        docx_path = os.path.join(default_storage.location, docx_filename)

        converter = Converter(full_pdf_path)
        converter.convert(docx_path, start=0, end=None)
        converter.close()

        # Send the DOCX file back as a download response
        with open(docx_path, 'rb') as docx_file:
            response = HttpResponse(docx_file, content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            response['Content-Disposition'] = f'attachment; filename="{docx_filename}"'
            return response

    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    finally:
        # Clean up temporary files
        if os.path.exists(full_pdf_path):
            os.remove(full_pdf_path)
        if os.path.exists(docx_path):
            os.remove(docx_path)

@api_view(['POST'])
def convert_pdf_to_text_file(request):

    if 'file' not in request.FILES:
        return Response({'error': 'No file uploaded'}, status=status.HTTP_400_BAD_REQUEST)

    pdf_file = request.FILES['file']
    if not pdf_file.name.endswith('.pdf'):
        return Response({'error': 'Invalid file format'}, status=status.HTTP_400_BAD_REQUEST)

    try:
        # Create a PDF reader object
        pdf_reader = PdfReader(pdf_file)
        text_buffer = io.BytesIO()

        for page in pdf_reader.pages:
            # Extract text from each page of the PDF
            
            text = page.extract_text()
            # output_file.write(text)
            text_buffer.write(text.encode('utf-8'))

        pdf_file.close()
       
        text_buffer.seek(0)

        # # Create a file response with the BytesIO object
        response = FileResponse(text_buffer, as_attachment=True, filename='converted.txt')

        return response;

    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['POST'])
def convert_pdf_to_ppt_file(request):

    if 'file' not in request.FILES:
        return Response({'error': 'No file uploaded'}, status=status.HTTP_400_BAD_REQUEST)

    pdf_file = request.FILES['file']
    if not pdf_file.name.endswith('.pdf'):
        return Response({'error': 'Invalid file format'}, status=status.HTTP_400_BAD_REQUEST)
    
    # Save the uploaded PDF file temporarily
    pdf_path = default_storage.save(pdf_file.name, pdf_file)
    full_pdf_path = os.path.join(default_storage.location, pdf_path)

    # Define the output PPTX file path
    pptx_filename = pdf_file.name.replace('.pdf', '.pptx')
    pptx_path = os.path.join(default_storage.location, pptx_filename)

    # Create a presentation object
    prs = Presentation()

    try:
        # Open the PDF using pdfplumber
        with pdfplumber.open(full_pdf_path) as pdf:
            for page_number, page in enumerate(pdf.pages):
                # Add a slide for each page
                slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
                slide = prs.slides.add_slide(slide_layout)

                # Extract text from the page
                text = page.extract_text()

                # Add text to the slide if available
                if text:
                    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                    text_frame = textbox.text_frame
                    text_frame.text = f"Page {page_number + 1}\n{text}"

                # Extract images from the page
                images = page.images
                for img_index, img in enumerate(images):
                    image_data = img['image']
                    image_path = f"{default_storage.location}/image_{page_number + 1}_{img_index + 1}.png"
                    with open(image_path, 'wb') as f:
                        f.write(image_data)

                    # Add image to the slide
                    slide.shapes.add_picture(image_path, Inches(1), Inches(1), Inches(5), Inches(5))

        # Save the PowerPoint presentation
        prs.save(pptx_path)

        # Send the PPTX file back as a download response
        with open(pptx_path, 'rb') as pptx_file:
            response = HttpResponse(pptx_file, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = f'attachment; filename="{pptx_filename}"'
            return response

    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)

    finally:
        if os.path.exists(full_pdf_path):
            os.remove(full_pdf_path)
        if os.path.exists(pptx_path):
            os.remove(pptx_path)

# @api_view(['POST'])
# def convert_pdf_to_excel_file(request):

#     if 'file' not in request.FILES:
#         return Response({'error': 'No file uploaded'}, status=status.HTTP_400_BAD_REQUEST)

#     pdf_file = request.FILES['file']
#     if not pdf_file.name.endswith('.pdf'):
#         return Response({'error': 'Invalid file format'}, status=status.HTTP_400_BAD_REQUEST)

#     try:
#        # Save the uploaded PDF file temporarily
#         pdf_path = default_storage.save(pdf_file.name, pdf_file)
#         full_pdf_path = os.path.join(default_storage.location, pdf_path)

#         # Define the output Excel (.xlsx) file path
#         excel_filename = pdf_file.name.replace('.pdf', '.xlsx')
#         excel_path = os.path.join(default_storage.location, excel_filename)

#         # Create an Excel workbook
#         workbook = Workbook()
#         sheet = workbook.active
#         sheet.title = "PDF Data"

#         try:
#             # Open the PDF using pdfplumber
#             with pdfplumber.open(full_pdf_path) as pdf:
#                 for page_number, page in enumerate(pdf.pages):
#                     # Extract tables from the page
#                     tables = page.extract_tables()

#                     # If tables are found, write them to the Excel sheet
#                     if tables:
#                         for table in tables:
#                             sheet.append([])  # Add an empty row to separate tables
#                             for row in table:
#                                 sheet.append(row)
#                     else:
#                         # If no tables, extract the text and put it in the Excel file
#                         text = page.extract_text()
#                         if text:
#                             sheet.append([f"Page {page_number + 1}"])
#                             for line in text.split("\n"):
#                                 sheet.append([line])

#             # Save the Excel workbook
#             workbook.save(excel_path)

#             # Send the Excel file back as a download response
#             with open(excel_path, 'rb') as excel_file:
#                 response = HttpResponse(excel_file, content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
#                 response['Content-Disposition'] = f'attachment; filename="{excel_filename}"'
#                 return response

#         except Exception as e:
#             return JsonResponse({"error": str(e)}, status=500)

#         finally:
#             if os.path.exists(full_pdf_path):
#                 os.remove(full_pdf_path)
#             if os.path.exists(excel_path):
#                 os.remove(excel_path)

#     except Exception as e:
#         return JsonResponse({"error": str(e)}, status=500)


@api_view(['POST'])
def merge_pdfs(request):
    # Check if files are uploaded
    if 'files' not in request.FILES or len(request.FILES.getlist('files')) < 2:
        return Response({'error': 'At least two PDF files are required'}, status=status.HTTP_400_BAD_REQUEST)
    
    pdf_files = request.FILES.getlist('files')
    
    # Validate all files to ensure they are PDFs
    for file in pdf_files:
        if not file.name.endswith('.pdf'):
            return Response({'error': f'Invalid file format: {file.name}'}, status=status.HTTP_400_BAD_REQUEST)
    
    try:
        # Create a PDF writer to store the merged content
        pdf_writer = PdfWriter()

        # Loop through each uploaded PDF and add its pages to the writer
        for pdf_file in pdf_files:
            pdf_reader = PdfReader(pdf_file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                pdf_writer.add_page(page)

        # Create a BytesIO object to store the merged PDF
        merged_pdf_buffer = io.BytesIO()

        # Write the merged PDF to the BytesIO object
        pdf_writer.write(merged_pdf_buffer)

        # Set the BytesIO object's position to the beginning
        merged_pdf_buffer.seek(0)

        # Create a file response with the BytesIO object
        response = FileResponse(merged_pdf_buffer, as_attachment=True, filename='merged.pdf')

        return response

    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    finally:
        for pdf_file in pdf_files:
            temp_pdf_path = os.path.join(default_storage.location, pdf_file.name)
            if os.path.exists(temp_pdf_path):
                os.remove(temp_pdf_path)


@api_view(['POST'])
def convert_image_to_pdf(request):
    if request.method == 'POST':
        if 'files' not in request.FILES:
            return JsonResponse({"error": "No image files uploaded"}, status=400)

        # Get the uploaded image files
        image_files = request.FILES.getlist('files')

        if not image_files:
            return JsonResponse({"error": "No image files uploaded"}, status=400)

        # Save the uploaded image files temporarily
        image_paths = []
        for image_file in image_files:
            image_path = default_storage.save(image_file.name, image_file)
            full_image_path = os.path.join(default_storage.location, image_path)
            image_paths.append(full_image_path)

        # Define the output PDF file path
        pdf_filename = "images_combined.pdf"
        pdf_path = os.path.join(default_storage.location, pdf_filename)

        try:
            # Create a PDF from the images
            images = []
            for img_path in image_paths:
                with Image.open(img_path) as img:
                    img = img.convert('RGB')  # Ensure the image is in RGB format
                    images.append(img)

            if images:
                # Save the first image as the PDF and append the rest
                images[0].save(pdf_path, save_all=True, append_images=images[1:], resolution=100.0)

            # Send the PDF file back as a download response
            with open(pdf_path, 'rb') as pdf_file:
                response = HttpResponse(pdf_file, content_type='application/pdf')
                response['Content-Disposition'] = f'attachment; filename="{pdf_filename}"'
                return response

        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

        finally:
            for image_path in image_paths:
                if os.path.exists(image_path):
                    os.remove(image_path)
            if os.path.exists(pdf_path):
                os.remove(pdf_path)

    return JsonResponse({"error": "Invalid request method"}, status=405)


@api_view(['POST'])
def convert_pdf_to_image(request):
    if request.method == 'POST':
        if 'file' not in request.FILES:
            return JsonResponse({"error": "No PDF file uploaded"}, status=400)

        # Get the uploaded PDF file
        pdf_file = request.FILES['file']

        # Save the uploaded PDF file temporarily
        pdf_path = default_storage.save(pdf_file.name, pdf_file)
        full_pdf_path = os.path.join(default_storage.location, pdf_path)

        try:
            # Convert PDF to image using pdf2image
            images = convert_from_path(full_pdf_path, 300)  # 300 DPI for high quality

            # Define the output image file path (assuming converting only the first page)
            image_filename = pdf_file.name.replace('.pdf', '.png')  # You can also use JPG or others
            image_path = os.path.join(default_storage.location, image_filename)

            # Save the first page as an image (PNG format)
            if images:
                images[0].save(image_path, 'PNG')

            # Send the image file back as a download response
            with open(image_path, 'rb') as image_file:
                response = HttpResponse(image_file, content_type='image/png')
                response['Content-Disposition'] = f'attachment; filename="{image_filename}"'
                return response

        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
        
        finally:
            if os.path.exists(full_pdf_path):
                os.remove(full_pdf_path)
            if os.path.exists(image_path):
                os.remove(image_path)
        
    return JsonResponse({"error": "Invalid request method"}, status=405)

@api_view(['POST'])
def convert_jpg_to_png(request):
    if request.method == 'POST':
        if 'file' not in request.FILES:
            return JsonResponse({"error": "No image file uploaded"}, status=400)

        # Get the uploaded image file
        image_file = request.FILES['file']

        # Validate if the uploaded file is a JPG image
        if not image_file.name.lower().endswith(('.jpg', '.jpeg')):
            return JsonResponse({"error": "Uploaded file is not a JPG image"}, status=400)

        # Save the uploaded image file temporarily
        image_path = default_storage.save(image_file.name, image_file)
        full_image_path = os.path.join(default_storage.location, image_path)

        # Define the output PNG file path
        png_filename = image_file.name.rsplit('.', 1)[0] + '.png'
        png_path = os.path.join(default_storage.location, png_filename)

        try:
            # Open the image using Pillow and convert to PNG
            with Image.open(full_image_path) as img:
                img = img.convert('RGB')  # Ensure the image is in RGB format
                img.save(png_path, "PNG")

            # Send the PNG file back as a download response
            with open(png_path, 'rb') as png_file:
                response = HttpResponse(png_file, content_type='image/png')
                response['Content-Disposition'] = f'attachment; filename="{png_filename}"'
                return response

        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

        finally:
            if os.path.exists(full_image_path):
                os.remove(full_image_path)
            if os.path.exists(png_path):
                os.remove(png_path)
    return JsonResponse({"error": "Invalid request method"}, status=405)


@api_view(['POST'])
def convert_png_to_jpg(request):
    if request.method == 'POST':
        if 'file' not in request.FILES:
            return JsonResponse({"error": "No image file uploaded"}, status=400)

        # Get the uploaded image file
        image_file = request.FILES['file']

        # Validate if the uploaded file is a PNG image
        if not image_file.name.lower().endswith('.png'):
            return JsonResponse({"error": "Uploaded file is not a PNG image"}, status=400)

        # Save the uploaded image file temporarily
        image_path = default_storage.save(image_file.name, image_file)
        full_image_path = os.path.join(default_storage.location, image_path)

        # Define the output JPG file path
        jpg_filename = image_file.name.rsplit('.', 1)[0] + '.jpg'
        jpg_path = os.path.join(default_storage.location, jpg_filename)

        try:
            # Open the image using Pillow and convert to JPG
            with Image.open(full_image_path) as img:
                img = img.convert('RGB')  # PNG images can have alpha channels, so we ensure RGB format
                img.save(jpg_path, "JPEG")

            # Send the JPG file back as a download response
            with open(jpg_path, 'rb') as jpg_file:
                response = HttpResponse(jpg_file, content_type='image/jpeg')
                response['Content-Disposition'] = f'attachment; filename="{jpg_filename}"'
                return response

        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

        finally:
            if os.path.exists(full_image_path):
                os.remove(full_image_path)
            if os.path.exists(jpg_path):
                os.remove(jpg_path)
    return JsonResponse({"error": "Invalid request method"}, status=405)

# @api_view(['POST'])
# def convert_svg_to_png(request):
#     if request.method == 'POST':
#         if 'file' not in request.FILES:
#             return JsonResponse({"error": "No SVG file uploaded"}, status=400)

#         # Get the uploaded SVG file
#         svg_file = request.FILES['file']

#         # Validate if the uploaded file is an SVG image
#         if not svg_file.name.lower().endswith('.svg'):
#             return JsonResponse({"error": "Uploaded file is not an SVG image"}, status=400)

#         # Save the uploaded SVG file temporarily
#         svg_path = default_storage.save(svg_file.name, svg_file)
#         full_svg_path = os.path.join(default_storage.location, svg_path)

#         # Define the output PNG file path
#         png_filename = svg_file.name.rsplit('.', 1)[0] + '.png'
#         png_path = os.path.join(default_storage.location, png_filename)

#         try:
#             # Convert the SVG file to PNG using CairoSVG
#             with open(full_svg_path, 'rb') as svg_input:
#                 svg_content = svg_input.read()
#                 cairosvg.svg2png(bytestring=svg_content, write_to=png_path)

#             # Send the PNG file back as a download response
#             with open(png_path, 'rb') as png_file:
#                 response = HttpResponse(png_file, content_type='image/png')
#                 response['Content-Disposition'] = f'attachment; filename="{png_filename}"'
#                 return response

#         except Exception as e:
#             return JsonResponse({"error": str(e)}, status=500)

#         finally:
#             if os.path.exists(full_svg_path):
#                 os.remove(full_svg_path)
#             if os.path.exists(png_path):
#                 os.remove(png_path)
            
#     return JsonResponse({"error": "Invalid request method"}, status=405)

# @api_view(['POST'])
# def convert_png_to_svg(request):
#     if request.method == 'POST':
#         if 'file' not in request.FILES:
#             return JsonResponse({"error": "No PNG file uploaded"}, status=400)

#         # Get the uploaded PNG image file
#         png_file = request.FILES['file']

#         # Validate if the uploaded file is a PNG image
#         if not png_file.name.lower().endswith('.png'):
#             return JsonResponse({"error": "Uploaded file is not a PNG image"}, status=400)

#         # Save the uploaded PNG image file temporarily
#         png_path = default_storage.save(png_file.name, png_file)
#         full_png_path = os.path.join(default_storage.location, png_path)

#         # Define the output PBM and SVG file paths
#         pbm_filename = png_file.name.rsplit('.', 1)[0] + '.pbm'
#         pbm_path = os.path.join(default_storage.location, pbm_filename)
#         svg_filename = png_file.name.rsplit('.', 1)[0] + '.svg'
#         svg_path = os.path.join(default_storage.location, svg_filename)

#         try:
#             # Open the PNG image using Pillow and convert to binary (1-bit pixels, black and white)
#             with Image.open(full_png_path) as img:
#                 img = img.convert('1')  # Convert to 1-bit pixels (black and white)
#                 img.save(pbm_path)  # Save as PBM (format determined by extension)

#             # Use Potrace to convert PBM to SVG
#             subprocess.run(['potrace', pbm_path, '-s', '-o', svg_path], check=True)

#             # Send the SVG file back as a download response
#             with open(svg_path, 'rb') as svg_output:
#                 response = HttpResponse(svg_output, content_type='image/svg+xml')
#                 response['Content-Disposition'] = f'attachment; filename="{svg_filename}"'
#                 return response

#         except subprocess.CalledProcessError as e:
#             return JsonResponse({"error": f"Potrace error: {e}"}, status=500)

#         except Exception as e:
#             return JsonResponse({"error": str(e)}, status=500)

#         finally:
#             if os.path.exists(full_png_path):
#                 os.remove(full_png_path)
#             if os.path.exists(svg_path):
#                 os.remove(svg_path)

#     return JsonResponse({"error": "Invalid request method"}, status=405)

@api_view(['POST'])
def compress_image(request):
    if request.method == 'POST':
        if 'file' not in request.FILES:
            return JsonResponse({"error": "No image file uploaded"}, status=400)

        # Get the uploaded image file
        image_file = request.FILES['file']

        # Validate if the uploaded file is an image
        if not image_file.name.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            return JsonResponse({"error": "Uploaded file is not a valid image"}, status=400)

        # Save the uploaded image file temporarily
        image_path = default_storage.save(image_file.name, image_file)
        full_image_path = os.path.join(default_storage.location, image_path)

        # Define the output compressed image file path
        compressed_filename = f"compressed_{image_file.name}"
        compressed_image_path = os.path.join(default_storage.location, compressed_filename)

        try:
            # Open the image using Pillow
            with Image.open(full_image_path) as img:
                # Resize the image (optional)
                new_size = (int(img.width * 0.6), int(img.height * 0.6))  # Reduce dimensions by 20%
                img = img.resize(new_size, Image.LANCZOS)

                img.save(compressed_image_path, format='jpeg', quality=60)  # Adjust quality as needed
                
            # Send the compressed image back as a download response
            with open(compressed_image_path, 'rb') as compressed_image_file:
                response = HttpResponse(compressed_image_file, content_type='image/jpeg')
                response['Content-Disposition'] = f'attachment; filename="{compressed_filename}"'
                return response

        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

        finally:
            if os.path.exists(full_image_path):
                os.remove(full_image_path)
            if os.path.exists(compressed_image_path):
                os.remove(compressed_image_path)

    return JsonResponse({"error": "Invalid request method"}, status=405)


@api_view(['POST'])
def compress_pdf(request):
    if 'file' not in request.FILES:
        return Response({'error': 'No file uploaded'}, status=status.HTTP_400_BAD_REQUEST)

    pdf_file = request.FILES['file']
    if not pdf_file.name.endswith('.pdf'):
        return Response({'error': 'Invalid file format, PDF required'}, status=status.HTTP_400_BAD_REQUEST)

    # Save the uploaded PDF file temporarily
    pdf_path = default_storage.save(pdf_file.name, pdf_file)
    full_pdf_path = os.path.join(default_storage.location, pdf_path)

    # Define the output path for the compressed PDF
    compressed_pdf_filename = "compressed_" + pdf_file.name
    compressed_pdf_path = os.path.join(default_storage.location, compressed_pdf_filename)

    try:
        # Open the PDF with PyMuPDF and compress it
        pdf_document = fitz.open(full_pdf_path)
        pdf_document.save(compressed_pdf_path, garbage=2, deflate=True, clean=True)
        pdf_document.close()

        # Send the compressed PDF file as a download response
        with open(compressed_pdf_path, 'rb') as compressed_pdf:
            response = HttpResponse(compressed_pdf, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{compressed_pdf_filename}"'
            return response

    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)

    finally:
        # Clean up the files
        if os.path.exists(full_pdf_path):
            os.remove(full_pdf_path)
        if os.path.exists(compressed_pdf_path):
            os.remove(compressed_pdf_path)


@api_view(['POST'])
def enhance_pdf_quality(request):
    if 'file' not in request.FILES:
        return Response({'error': 'No file uploaded'}, status=status.HTTP_400_BAD_REQUEST)

    pdf_file = request.FILES['file']
    if not pdf_file.name.endswith('.pdf'):
        return Response({'error': 'Invalid file format, PDF required'}, status=status.HTTP_400_BAD_REQUEST)

    # Save the uploaded PDF file temporarily
    pdf_path = default_storage.save(pdf_file.name, pdf_file)
    full_pdf_path = os.path.join(default_storage.location, pdf_path)

    # Define the output path for the compressed PDF
    compressed_pdf_filename = "compressed_" + pdf_file.name
    compressed_pdf_path = os.path.join(default_storage.location, compressed_pdf_filename)

    try:
        # Open the PDF with PyMuPDF
        pdf_document = fitz.open(full_pdf_path)

        # Loop through each page to reduce image size and quality
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            images = page.get_images(full=True)

            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]

                # Use PIL to open the image
                with Image.open(BytesIO(image_bytes)) as image:
                    # Convert image to RGB if it is in a different mode
                    if image.mode != "RGB":
                        image = image.convert("RGB")

                    # Resize the image to 50% of its original dimensions
                    new_size = (int(image.width * 1.5), int(image.height * 1.5))

                    resized_image = image.resize(new_size, Image.LANCZOS)

                    # Save the resized image to a BytesIO object with JPEG format and reduced quality
                    img_byte_arr = BytesIO()
                    resized_image.save(img_byte_arr, format="jpeg", quality=150)

                    img_byte_arr = img_byte_arr.getvalue()

                    # Replace the original image with the resized and compressed image in the PDF
                    page.insert_image(page.rect, stream=img_byte_arr, keep_proportion=True)

        # Save the compressed PDF
        pdf_document.save(compressed_pdf_path, garbage=4, deflate=True, clean=True)
        pdf_document.close()

        # Send the compressed PDF file as a download response
        with open(compressed_pdf_path, 'rb') as compressed_pdf:
            response = HttpResponse(compressed_pdf, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{compressed_pdf_filename}"'
            return response

    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)

    finally:
        # Clean up the files
        if os.path.exists(full_pdf_path):
            os.remove(full_pdf_path)
        if os.path.exists(compressed_pdf_path):
            os.remove(compressed_pdf_path)


# @api_view(['POST'])
# def convert_docx_to_pdf(request):
#     if 'file' not in request.FILES:
#         return Response({'error': 'No file uploaded'}, status=status.HTTP_400_BAD_REQUEST)

#     uploaded_file = request.FILES['file']
#     file_name, file_extension = os.path.splitext(uploaded_file.name)

#     # Check for .docx file type
#     if file_extension.lower() != '.docx':
#         return Response({'error': 'Invalid file format. Only .docx files are allowed.'}, status=status.HTTP_400_BAD_REQUEST)

#     # Save the uploaded file temporarily
#     file_path = default_storage.save(uploaded_file.name, uploaded_file)
#     full_file_path = os.path.join(default_storage.location, file_path)

#     # Define output path for the PDF
#     pdf_filename = f"{file_name}.pdf"
#     pdf_path = os.path.join(default_storage.location, pdf_filename)

#     try:
#         # Convert the .docx file to PDF
#         convert_to_pdf(full_file_path, pdf_path)

#         # Send the converted PDF file as a download response
#         with open(pdf_path, 'rb') as pdf_file:
#             response = HttpResponse(pdf_file, content_type='application/pdf')
#             response['Content-Disposition'] = f'attachment; filename="{pdf_filename}"'
#             return response

#     except Exception as e:
#         return JsonResponse({"error": str(e)}, status=500)

#     finally:
#         # Clean up the files
#         if os.path.exists(full_file_path):
#             os.remove(full_file_path)
#         if os.path.exists(pdf_path):
#             os.remove(pdf_path)


# def convert_to_pdf(docx_path, pdf_path):
#     # Open the .docx file and initialize the PDF
#     doc = DocxDocument(docx_path)
#     pdf = FPDF()
#     pdf.set_auto_page_break(auto=True, margin=15)
#     pdf.add_page()
#     pdf.set_font("Arial", size=12)

#     # Add each paragraph from the .docx to the PDF
#     for paragraph in doc.paragraphs:
#         pdf.multi_cell(0, 10, paragraph.text)

#     # Save the PDF
#     pdf.output(pdf_path)